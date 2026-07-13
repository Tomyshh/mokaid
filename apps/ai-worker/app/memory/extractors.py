"""Central text extraction for every document format Mokaid supports.

One entry point — ``extract_bytes(data, filename, mime_type)`` — used both by
the knowledge ingestion pipeline (so PDF/DOCX/XLSX/PPTX files become part of
the agents' vector knowledge base) and by the ``extract_document_text``
mission tool. Each format has a dedicated parser with graceful fallbacks;
binary formats we cannot parse return ``None`` so callers can decide to OCR
with vision or fail cleanly.
"""

import io
from dataclasses import dataclass, field

import structlog

log = structlog.get_logger()

# Formats we can turn into text without an LLM. Keep in sync with the
# Phoenix @indexable_extensions list (agent feed-data + knowledge ingestion).
EXTRACTABLE_EXTENSIONS = (
    ".pdf",
    ".docx",
    ".xlsx",
    ".xlsm",
    ".xls",
    ".pptx",
    ".rtf",
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".tsv",
    ".json",
    ".html",
    ".htm",
    ".xml",
    ".yaml",
    ".yml",
)

_TEXT_EXTENSIONS = (
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".tsv",
    ".json",
    ".html",
    ".htm",
    ".xml",
    ".yaml",
    ".yml",
)

# Hard caps so a huge workbook or deck can't blow memory or the embedding
# budget. Ingestion chunks downstream, so generous but bounded.
_MAX_CHARS = 800_000
_MAX_SHEET_ROWS = 2_000


@dataclass
class ExtractResult:
    text: str
    format: str
    metadata: dict = field(default_factory=dict)


def _ext(filename: str) -> str:
    name = (filename or "").lower()
    dot = name.rfind(".")
    return name[dot:] if dot >= 0 else ""


def looks_like_text(value: str) -> bool:
    """True when the string is mostly clean text (not decoded binary).

    A raw PDF/binary decoded as UTF-8 is dominated by control chars and the
    U+FFFD replacement character — we must never return that as "extracted
    text": it pollutes outputs and PostgreSQL rejects NUL bytes.
    """
    if not value:
        return False
    sample = value[:4000]
    if sample.count("\ufffd") / len(sample) > 0.1:
        return False
    clean = sum(1 for ch in sample if (ch.isprintable() and ch != "\ufffd") or ch in "\t\n\r")
    return clean / len(sample) >= 0.85


def is_extractable(filename: str, mime_type: str | None = None) -> bool:
    if _ext(filename) in EXTRACTABLE_EXTENSIONS:
        return True
    mime = (mime_type or "").lower()
    return mime.startswith("text/") or mime in (
        "application/pdf",
        "application/json",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-excel",
        "application/rtf",
    )


def extract_bytes(
    data: bytes, filename: str = "", mime_type: str | None = None
) -> ExtractResult | None:
    """Extract readable text from raw file bytes. Returns None when the
    format is binary and unparseable (caller may fall back to vision OCR)."""
    if not data:
        return None

    ext = _ext(filename)
    mime = (mime_type or "").lower()

    if ext == ".pdf" or data[:5] == b"%PDF-" or mime == "application/pdf":
        return _extract_pdf(data)
    if ext == ".docx" or mime.endswith("wordprocessingml.document"):
        return _extract_docx(data)
    if ext in (".xlsx", ".xlsm") or mime.endswith("spreadsheetml.sheet"):
        return _extract_xlsx(data)
    if ext == ".xls" or mime == "application/vnd.ms-excel":
        return _extract_xls(data)
    if ext == ".pptx" or mime.endswith("presentationml.presentation"):
        return _extract_pptx(data)
    if ext == ".rtf" or mime == "application/rtf":
        return _extract_rtf(data)

    if ext in _TEXT_EXTENSIONS or mime.startswith("text/") or mime == "application/json":
        return _extract_plain_text(data)

    # Unknown extension: try plain text decode as a last cheap attempt.
    return _extract_plain_text(data)


# ---------- PDF ----------


def _extract_pdf(data: bytes) -> ExtractResult | None:
    text = ""
    page_count = 0

    # 1) PyMuPDF — best quality (layout-aware, fast).
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(stream=data, filetype="pdf")
        page_count = doc.page_count
        text = "\n\n".join(page.get_text() for page in doc)
        doc.close()
    except ImportError:
        text = ""
    except Exception as exc:
        log.warning("pdf_fitz_failed", error=str(exc))
        text = ""

    # 2) pypdf fallback (pure-python).
    if not text.strip():
        try:
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(data))
            page_count = len(reader.pages)
            text = "\n\n".join((page.extract_text() or "") for page in reader.pages)
        except Exception as exc:
            log.warning("pdf_pypdf_failed", error=str(exc))
            text = ""

    # 3) Annotations (FreeText overlays, form fields, signatures) live
    # outside the page content stream — extract them separately.
    annot_text = extract_pdf_annotations(data)
    if annot_text:
        text = (text + "\n\n--- Annotations & Signatures ---\n" + annot_text).strip()

    text = text.strip()
    if not text:
        return None
    return ExtractResult(text[:_MAX_CHARS], "pdf", {"page_count": page_count})


def extract_pdf_annotations(doc_bytes: bytes) -> str:
    """Extract text from PDF annotations (FreeText, stamps, signatures)."""
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(doc_bytes))
    except Exception:
        return ""

    lines: list[str] = []
    for page_num, page in enumerate(reader.pages, 1):
        if "/Annots" not in page:
            continue
        try:
            annots_ref = page["/Annots"]
            annots = annots_ref.get_object() if hasattr(annots_ref, "get_object") else annots_ref
            if not isinstance(annots, list):
                continue
        except Exception:
            continue

        for a in annots:
            try:
                ao = a.get_object() if hasattr(a, "get_object") else a
                subtype = str(ao.get("/Subtype", ""))
                contents = ao.get("/Contents", "")

                if subtype == "/FreeText" and contents:
                    lines.append(f"[Page {page_num} annotation] {contents}")

                if str(ao.get("/FT", "")) == "/Sig":
                    v = ao.get("/V")
                    if v:
                        sig = v.get_object() if hasattr(v, "get_object") else v
                        name = sig.get("/Name", "")
                        reason = sig.get("/Reason", "")
                        if name:
                            lines.append(
                                f"[Page {page_num} digital signature] Signer: {name}"
                                + (f", Reason: {reason}" if reason else "")
                            )
            except Exception:
                continue

    return "\n".join(lines)


# ---------- Word ----------


def _extract_docx(data: bytes) -> ExtractResult | None:
    try:
        import docx
    except ImportError:
        log.warning("docx_parser_missing")
        return None

    try:
        document = docx.Document(io.BytesIO(data))
    except Exception as exc:
        log.warning("docx_parse_failed", error=str(exc))
        return None

    parts: list[str] = []
    for para in document.paragraphs:
        content = para.text.strip()
        if not content:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        # Preserve heading structure as markdown so the structure-aware
        # splitter can chunk along sections.
        if style.startswith("heading"):
            try:
                level = int(style.rsplit(" ", 1)[-1])
            except ValueError:
                level = 2
            parts.append("#" * min(level, 6) + " " + content)
        else:
            parts.append(content)

    for table in document.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            header_sep = "| " + " | ".join("---" for _ in table.rows[0].cells) + " |"
            rows.insert(1, header_sep)
            parts.append("\n".join(rows))

    text = "\n\n".join(parts).strip()
    if not text:
        return None
    return ExtractResult(
        text[:_MAX_CHARS],
        "docx",
        {"paragraph_count": len(document.paragraphs), "table_count": len(document.tables)},
    )


# ---------- Excel ----------


def _cell_str(value) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value == int(value):
        return str(int(value))
    return str(value).replace("\n", " ").strip()


def _rows_to_markdown(rows: list[list[str]]) -> str:
    """Render rows as a markdown table, skipping fully empty rows."""
    filled = [r for r in rows if any(c for c in r)]
    if not filled:
        return ""
    width = max(len(r) for r in filled)
    normalized = [r + [""] * (width - len(r)) for r in filled]
    lines = ["| " + " | ".join(normalized[0]) + " |"]
    lines.append("| " + " | ".join("---" for _ in range(width)) + " |")
    for row in normalized[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _extract_xlsx(data: bytes) -> ExtractResult | None:
    try:
        import openpyxl
    except ImportError:
        log.warning("xlsx_parser_missing")
        return None

    try:
        workbook = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:
        log.warning("xlsx_parse_failed", error=str(exc))
        return None

    parts: list[str] = []
    sheet_names: list[str] = []
    for sheet in workbook.worksheets:
        sheet_names.append(sheet.title)
        rows: list[list[str]] = []
        truncated = False
        for i, row in enumerate(sheet.iter_rows(values_only=True)):
            if i >= _MAX_SHEET_ROWS:
                truncated = True
                break
            rows.append([_cell_str(v) for v in row])
        table = _rows_to_markdown(rows)
        if table:
            section = f"## Sheet: {sheet.title}\n\n{table}"
            if truncated:
                section += f"\n\n(truncated at {_MAX_SHEET_ROWS} rows)"
            parts.append(section)
    workbook.close()

    text = "\n\n".join(parts).strip()
    if not text:
        return None
    return ExtractResult(text[:_MAX_CHARS], "xlsx", {"sheets": sheet_names})


def _extract_xls(data: bytes) -> ExtractResult | None:
    try:
        import xlrd
    except ImportError:
        log.warning("xls_parser_missing")
        return None

    try:
        workbook = xlrd.open_workbook(file_contents=data)
    except Exception as exc:
        log.warning("xls_parse_failed", error=str(exc))
        return None

    parts: list[str] = []
    sheet_names: list[str] = []
    for sheet in workbook.sheets():
        sheet_names.append(sheet.name)
        rows: list[list[str]] = []
        limit = min(sheet.nrows, _MAX_SHEET_ROWS)
        for i in range(limit):
            rows.append([_cell_str(v) for v in sheet.row_values(i)])
        table = _rows_to_markdown(rows)
        if table:
            section = f"## Sheet: {sheet.name}\n\n{table}"
            if sheet.nrows > _MAX_SHEET_ROWS:
                section += f"\n\n(truncated at {_MAX_SHEET_ROWS} rows)"
            parts.append(section)

    text = "\n\n".join(parts).strip()
    if not text:
        return None
    return ExtractResult(text[:_MAX_CHARS], "xls", {"sheets": sheet_names})


# ---------- PowerPoint ----------


def _extract_pptx(data: bytes) -> ExtractResult | None:
    try:
        from pptx import Presentation
    except ImportError:
        log.warning("pptx_parser_missing")
        return None

    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:
        log.warning("pptx_parse_failed", error=str(exc))
        return None

    parts: list[str] = []
    for num, slide in enumerate(presentation.slides, 1):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    content = "".join(run.text for run in para.runs).strip()
                    if content:
                        lines.append(content)
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    rows.append(cells)
                table = _rows_to_markdown(rows)
                if table:
                    lines.append(table)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        section = f"## Slide {num}\n\n" + "\n".join(lines)
        if notes:
            section += f"\n\nSpeaker notes: {notes}"
        parts.append(section)

    text = "\n\n".join(parts).strip()
    if not text:
        return None
    return ExtractResult(text[:_MAX_CHARS], "pptx", {"slide_count": len(presentation.slides)})


# ---------- RTF / plain text ----------


def _extract_rtf(data: bytes) -> ExtractResult | None:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        log.warning("rtf_parser_missing")
        return None

    try:
        text = rtf_to_text(data.decode("utf-8", errors="replace")).strip()
    except Exception as exc:
        log.warning("rtf_parse_failed", error=str(exc))
        return None

    if not text or not looks_like_text(text):
        return None
    return ExtractResult(text[:_MAX_CHARS], "rtf", {})


def _extract_plain_text(data: bytes) -> ExtractResult | None:
    decoded = data.decode("utf-8", errors="replace").strip()
    if not decoded or not looks_like_text(decoded):
        return None
    return ExtractResult(decoded[:_MAX_CHARS], "text", {})
