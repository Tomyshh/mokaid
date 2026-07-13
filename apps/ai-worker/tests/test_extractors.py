"""Extractor tests with generated fixtures — no LLM, no network.

Each fixture is built in-memory with the same libraries the production
parsers rely on (python-docx, openpyxl, python-pptx, PyMuPDF), so the tests
exercise the real parsing paths for every supported format.
"""

import io

import pytest

from app.memory import extractors
from app.memory.extractors import ExtractResult, extract_bytes, is_extractable, looks_like_text

# ---------- fixture builders ----------


def make_pdf(texts: list[str]) -> bytes:
    import fitz

    doc = fitz.open()
    for text in texts:
        page = doc.new_page()
        page.insert_text((72, 72), text)
    data = doc.tobytes()
    doc.close()
    return data


def make_docx() -> bytes:
    import docx

    document = docx.Document()
    document.add_heading("Contrat de prestation", level=1)
    document.add_paragraph("Le prestataire s'engage à livrer le rapport avant le 30 juin.")
    document.add_heading("Clauses financières", level=2)
    document.add_paragraph("Le montant total est de 12 000 euros hors taxes.")
    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Poste"
    table.rows[0].cells[1].text = "Montant"
    table.rows[1].cells[0].text = "Audit"
    table.rows[1].cells[1].text = "5000"
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def make_xlsx() -> bytes:
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Comptes Q3"
    sheet.append(["Poste", "Montant"])
    sheet.append(["Revenus", 150000])
    sheet.append(["Charges", 90000])
    other = workbook.create_sheet("Prévisions")
    other.append(["Trimestre", "Objectif"])
    other.append(["Q4", 200000])
    buf = io.BytesIO()
    workbook.save(buf)
    return buf.getvalue()


def make_pptx() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Stratégie SEO 2026"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
    box.text_frame.text = "Objectif : doubler le trafic organique."
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


RTF_SAMPLE = b"{\\rtf1\\ansi Le patient pr\\'e9sente une am\\'e9lioration notable.}"


# ---------- per-format extraction ----------


def test_pdf_extraction():
    data = make_pdf(["Diagnostic: hypertension stade 2.", "Traitement recommandé: IEC."])
    result = extract_bytes(data, "rapport.pdf")
    assert isinstance(result, ExtractResult)
    assert result.format == "pdf"
    assert result.metadata["page_count"] == 2
    assert "hypertension" in result.text
    assert "IEC" in result.text


def test_pdf_detected_by_magic_bytes_without_extension():
    data = make_pdf(["Contenu sans extension."])
    result = extract_bytes(data, "piece_jointe")
    assert result is not None
    assert result.format == "pdf"


def test_docx_extraction_headings_and_tables():
    result = extract_bytes(make_docx(), "contrat.docx")
    assert result is not None
    assert result.format == "docx"
    # Headings become markdown so the splitter can chunk along sections.
    assert "# Contrat de prestation" in result.text
    assert "## Clauses financières" in result.text
    assert "12 000 euros" in result.text
    # Tables become markdown tables.
    assert "| Audit | 5000 |" in result.text


def test_xlsx_extraction_sheets_as_markdown():
    result = extract_bytes(make_xlsx(), "comptes.xlsx")
    assert result is not None
    assert result.format == "xlsx"
    assert result.metadata["sheets"] == ["Comptes Q3", "Prévisions"]
    assert "## Sheet: Comptes Q3" in result.text
    assert "| Revenus | 150000 |" in result.text
    assert "| Q4 | 200000 |" in result.text


def test_pptx_extraction_slides():
    result = extract_bytes(make_pptx(), "deck.pptx")
    assert result is not None
    assert result.format == "pptx"
    assert "## Slide 1" in result.text
    assert "Stratégie SEO 2026" in result.text
    assert "trafic organique" in result.text


def test_rtf_extraction():
    result = extract_bytes(RTF_SAMPLE, "note.rtf")
    assert result is not None
    assert result.format == "rtf"
    assert "amélioration notable" in result.text


def test_plain_text_and_csv():
    txt = extract_bytes(b"Simple note de brief SEO.", "brief.txt")
    assert txt is not None and txt.format == "text"

    csv = extract_bytes(b"mot,volume\nseo,1000\n", "keywords.csv")
    assert csv is not None
    assert "seo,1000" in csv.text


def test_unknown_binary_returns_none():
    # PNG magic + random bytes: must never come back as "text".
    data = b"\x89PNG\r\n\x1a\n" + bytes(range(256)) * 8
    assert extract_bytes(data, "image.png") is None


def test_empty_input_returns_none():
    assert extract_bytes(b"", "vide.pdf") is None


def test_corrupt_docx_returns_none():
    assert extract_bytes(b"not a zip archive at all", "broken.docx") is None


# ---------- guards & helpers ----------


def test_looks_like_text_rejects_decoded_binary():
    binary_decoded = ("\ufffd" * 50 + "abc") * 20
    assert not looks_like_text(binary_decoded)
    assert looks_like_text("Un paragraphe parfaitement lisible.\nAvec des lignes.")


@pytest.mark.parametrize(
    ("filename", "mime", "expected"),
    [
        ("doc.pdf", None, True),
        ("feuille.xlsx", None, True),
        ("pres.pptx", None, True),
        ("archive.zip", None, False),
        ("photo.jpg", "image/jpeg", False),
        ("sans_ext", "application/pdf", True),
        ("sans_ext", "text/plain", True),
    ],
)
def test_is_extractable(filename, mime, expected):
    assert is_extractable(filename, mime) is expected


def test_max_chars_cap(monkeypatch):
    monkeypatch.setattr(extractors, "_MAX_CHARS", 100)
    result = extract_bytes(("mot " * 200).encode(), "long.txt")
    assert result is not None
    assert len(result.text) == 100
